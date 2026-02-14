"""PowerPoint presentation renderer using python-pptx.

Supports multiple slide types including data visualization slides
(stat_callout, comparison, timeline, chart) for skill-enhanced output.
"""

import logging
from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from app.models.schemas import GeneratedPresentation

logger = logging.getLogger(__name__)


class PPTXRenderer:
    """Renders presentations as PPTX files with enhanced slide types."""

    # --- Color palettes ---
    # Primary palette
    PRIMARY_COLOR = RGBColor(37, 99, 235)    # #2563eb — blue
    SECONDARY_COLOR = RGBColor(30, 64, 175)  # #1e40af — dark blue
    ACCENT_COLOR = RGBColor(16, 185, 129)    # #10b981 — emerald
    ACCENT_WARM = RGBColor(245, 158, 11)     # #f59e0b — amber
    WHITE = RGBColor(255, 255, 255)
    DARK_TEXT = RGBColor(26, 26, 26)
    GRAY_TEXT = RGBColor(107, 114, 128)
    LIGHT_BG = RGBColor(243, 244, 246)       # #f3f4f6 — gray-100
    LIGHT_BLUE_BG = RGBColor(239, 246, 255)  # #eff6ff — blue-50

    # Chart colors (for bars/pie slices)
    CHART_COLORS = [
        RGBColor(37, 99, 235),    # blue
        RGBColor(16, 185, 129),   # emerald
        RGBColor(245, 158, 11),   # amber
        RGBColor(239, 68, 68),    # red
        RGBColor(139, 92, 246),   # purple
        RGBColor(14, 165, 233),   # sky
        RGBColor(249, 115, 22),   # orange
        RGBColor(34, 197, 94),    # green
    ]

    # Slide dimensions (16:9)
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    def render(self, presentation: GeneratedPresentation) -> BytesIO:
        """Render a presentation to PPTX."""
        try:
            prs = Presentation()
            prs.slide_width = self.SLIDE_WIDTH
            prs.slide_height = self.SLIDE_HEIGHT

            for slide_data in presentation.slides:
                slide_type = slide_data.get("type", "content")

                renderer_map = {
                    "title": self._add_title_slide,
                    "section": self._add_section_slide,
                    "content": self._add_content_slide,
                    "key_findings": self._add_findings_slide,
                    "stat_callout": self._add_stat_callout_slide,
                    "comparison": self._add_comparison_slide,
                    "timeline": self._add_timeline_slide,
                    "chart": self._add_chart_slide,
                    "recommendations": self._add_recommendations_slide,
                    "closing": self._add_closing_slide,
                }

                renderer = renderer_map.get(slide_type, self._add_content_slide)
                renderer(prs, slide_data)

            output = BytesIO()
            prs.save(output)
            output.seek(0)
            return output

        except Exception as e:
            logger.error(f"PPTX rendering failed: {e}")
            raise

    # ===========================================================
    # Helper methods
    # ===========================================================

    def _add_bg(self, slide, color: RGBColor) -> None:
        """Set solid background color on a slide."""
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_rect(self, slide, left, top, width, height, color: RGBColor):
        """Add a filled rectangle shape (no border)."""
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def _add_text(
        self, slide, left, top, width, height,
        text: str, font_size: int = 20,
        color: RGBColor = None, bold: bool = False,
        alignment=PP_ALIGN.LEFT, word_wrap: bool = True,
    ):
        """Add a text box with styling."""
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = word_wrap
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        if color:
            p.font.color.rgb = color
        p.alignment = alignment
        return box

    def _add_speaker_notes(self, slide, notes: str | None) -> None:
        """Add speaker notes to a slide if provided."""
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    # ===========================================================
    # Slide renderers
    # ===========================================================

    def _add_title_slide(self, prs: Presentation, data: dict) -> None:
        """Add a title slide with branded background."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_bg(slide, self.PRIMARY_COLOR)

        self._add_text(
            slide, Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5),
            text=data.get("title", "Presentation"),
            font_size=44, bold=True, color=self.WHITE,
            alignment=PP_ALIGN.CENTER,
        )

        if subtitle := data.get("subtitle"):
            self._add_text(
                slide, Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8),
                text=subtitle, font_size=24, color=self.WHITE,
                alignment=PP_ALIGN.CENTER,
            )

        self._add_speaker_notes(slide, data.get("notes"))

    def _add_section_slide(self, prs: Presentation, data: dict) -> None:
        """Add a section divider slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Left accent panel
        self._add_rect(
            slide, Inches(0), Inches(0),
            Inches(6.667), Inches(7.5),
            self.PRIMARY_COLOR,
        )

        self._add_text(
            slide, Inches(0.5), Inches(3), Inches(6), Inches(1.5),
            text=data.get("title", "Section"),
            font_size=36, bold=True, color=self.WHITE,
        )

        self._add_speaker_notes(slide, data.get("notes"))

    def _add_content_slide(self, prs: Presentation, data: dict) -> None:
        """Add a content slide with bullets."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title
        self._add_text(
            slide, Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8),
            text=data.get("title", ""),
            font_size=32, bold=True, color=self.PRIMARY_COLOR,
        )

        # Accent line
        self._add_rect(
            slide, Inches(0.5), Inches(1.1),
            Inches(2), Inches(0.05),
            self.PRIMARY_COLOR,
        )

        # Bullet points
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5)
        )
        tf = content_box.text_frame
        tf.word_wrap = True

        bullets = data.get("bullets", [])
        for i, bullet in enumerate(bullets[:6]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"\u2022 {bullet}"
            p.font.size = Pt(20)
            p.font.color.rgb = self.DARK_TEXT
            p.space_after = Pt(12)

        self._add_speaker_notes(slide, data.get("notes"))

    def _add_findings_slide(self, prs: Presentation, data: dict) -> None:
        """Add a key findings slide."""
        self._add_content_slide(prs, {
            "title": data.get("title", "Key Findings"),
            "bullets": data.get("findings", []),
            "notes": data.get("notes"),
        })

    def _add_recommendations_slide(self, prs: Presentation, data: dict) -> None:
        """Add a recommendations slide."""
        self._add_content_slide(prs, {
            "title": data.get("title", "Recommendations"),
            "bullets": data.get("items", []),
            "notes": data.get("notes"),
        })

    def _add_stat_callout_slide(self, prs: Presentation, data: dict) -> None:
        """Add a stat callout slide — one big metric with context.

        Layout:
        - Light background
        - Title at top
        - Large centered stat value
        - Context line below
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_bg(slide, self.LIGHT_BLUE_BG)

        # Title
        self._add_text(
            slide, Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8),
            text=data.get("title", ""),
            font_size=28, bold=True, color=self.PRIMARY_COLOR,
            alignment=PP_ALIGN.CENTER,
        )

        # Accent line centered
        self._add_rect(
            slide, Inches(5.667), Inches(1.5),
            Inches(2), Inches(0.05),
            self.PRIMARY_COLOR,
        )

        # Big stat value
        stat_value = data.get("stat_value", "—")
        self._add_text(
            slide, Inches(0.5), Inches(2.2), Inches(12.333), Inches(2.5),
            text=stat_value,
            font_size=72, bold=True, color=self.PRIMARY_COLOR,
            alignment=PP_ALIGN.CENTER,
        )

        # Context line
        stat_context = data.get("stat_context", "")
        if stat_context:
            self._add_text(
                slide, Inches(1.5), Inches(4.8), Inches(10.333), Inches(1),
                text=stat_context,
                font_size=22, color=self.GRAY_TEXT,
                alignment=PP_ALIGN.CENTER,
            )

        self._add_speaker_notes(slide, data.get("notes"))

    def _add_comparison_slide(self, prs: Presentation, data: dict) -> None:
        """Add a comparison slide — two columns side by side.

        Layout:
        - Title at top
        - Left column (label + items) | Right column (label + items)
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title
        self._add_text(
            slide, Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8),
            text=data.get("title", "Comparison"),
            font_size=32, bold=True, color=self.PRIMARY_COLOR,
            alignment=PP_ALIGN.CENTER,
        )

        # Accent line
        self._add_rect(
            slide, Inches(5.667), Inches(1.1),
            Inches(2), Inches(0.05),
            self.PRIMARY_COLOR,
        )

        # Divider line (vertical, centered)
        self._add_rect(
            slide, Inches(6.567), Inches(1.5),
            Inches(0.04), Inches(5.5),
            self.LIGHT_BG,
        )

        # Left column
        left_label = data.get("left_label", "Option A")
        self._add_text(
            slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.6),
            text=left_label,
            font_size=24, bold=True, color=self.PRIMARY_COLOR,
        )

        left_items = data.get("left_items", [])
        left_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.3), Inches(5.8), Inches(4.5)
        )
        tf_left = left_box.text_frame
        tf_left.word_wrap = True
        for i, item in enumerate(left_items[:6]):
            p = tf_left.paragraphs[0] if i == 0 else tf_left.add_paragraph()
            p.text = f"\u2022 {item}"
            p.font.size = Pt(18)
            p.font.color.rgb = self.DARK_TEXT
            p.space_after = Pt(10)

        # Right column
        right_label = data.get("right_label", "Option B")
        self._add_text(
            slide, Inches(7), Inches(1.5), Inches(5.8), Inches(0.6),
            text=right_label,
            font_size=24, bold=True, color=self.ACCENT_COLOR,
        )

        right_items = data.get("right_items", [])
        right_box = slide.shapes.add_textbox(
            Inches(7), Inches(2.3), Inches(5.8), Inches(4.5)
        )
        tf_right = right_box.text_frame
        tf_right.word_wrap = True
        for i, item in enumerate(right_items[:6]):
            p = tf_right.paragraphs[0] if i == 0 else tf_right.add_paragraph()
            p.text = f"\u2022 {item}"
            p.font.size = Pt(18)
            p.font.color.rgb = self.DARK_TEXT
            p.space_after = Pt(10)

        self._add_speaker_notes(slide, data.get("notes"))

    def _add_timeline_slide(self, prs: Presentation, data: dict) -> None:
        """Add a timeline slide with sequential events.

        Layout:
        - Title at top
        - Horizontal line with event markers
        - Date + description below each marker
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title
        self._add_text(
            slide, Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8),
            text=data.get("title", "Timeline"),
            font_size=32, bold=True, color=self.PRIMARY_COLOR,
        )

        events = data.get("events", [])
        if not events:
            self._add_speaker_notes(slide, data.get("notes"))
            return

        num_events = min(len(events), 6)
        events = events[:num_events]

        # Horizontal timeline bar
        bar_y = Inches(3.2)
        bar_left = Inches(1)
        bar_width = Inches(11.333)
        self._add_rect(slide, bar_left, bar_y, bar_width, Inches(0.06), self.PRIMARY_COLOR)

        # Event markers and text
        spacing = bar_width / (num_events + 1) if num_events > 0 else bar_width
        for idx, event in enumerate(events):
            cx = bar_left + spacing * (idx + 1)

            # Marker circle (small square with round appearance)
            marker = self._add_rect(
                slide,
                cx - Inches(0.15), bar_y - Inches(0.12),
                Inches(0.3), Inches(0.3),
                self.ACCENT_COLOR,
            )

            # Date/label above the bar
            date_text = event.get("date", event.get("label", f"Step {idx+1}"))
            self._add_text(
                slide,
                cx - Inches(0.8), bar_y - Inches(1.2),
                Inches(1.6), Inches(0.8),
                text=date_text,
                font_size=14, bold=True, color=self.PRIMARY_COLOR,
                alignment=PP_ALIGN.CENTER,
            )

            # Description below the bar
            desc = event.get("description", "")
            if desc:
                self._add_text(
                    slide,
                    cx - Inches(0.9), bar_y + Inches(0.5),
                    Inches(1.8), Inches(1.5),
                    text=desc,
                    font_size=12, color=self.DARK_TEXT,
                    alignment=PP_ALIGN.CENTER,
                )

        self._add_speaker_notes(slide, data.get("notes"))

    def _add_chart_slide(self, prs: Presentation, data: dict) -> None:
        """Add a chart slide with a matplotlib-generated chart image.

        Uses matplotlib to render the chart as a PNG and embeds it in the slide.
        Falls back to a text-based representation if matplotlib is unavailable.
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title
        chart_title = data.get("chart_title") or data.get("title", "Chart")
        self._add_text(
            slide, Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8),
            text=chart_title,
            font_size=28, bold=True, color=self.PRIMARY_COLOR,
            alignment=PP_ALIGN.CENTER,
        )

        data_labels = data.get("data_labels", [])
        data_values = data.get("data_values", [])
        chart_type = data.get("chart_type", "bar")

        if not data_labels or not data_values:
            # No data — show placeholder
            self._add_text(
                slide, Inches(2), Inches(3), Inches(9.333), Inches(1),
                text="[Chart data unavailable]",
                font_size=20, color=self.GRAY_TEXT,
                alignment=PP_ALIGN.CENTER,
            )
            self._add_speaker_notes(slide, data.get("notes"))
            return

        try:
            chart_image = self._render_chart_image(
                chart_type, data_labels, data_values, chart_title
            )
            # Add chart image to slide (centered)
            slide.shapes.add_picture(
                chart_image,
                Inches(1.5), Inches(1.3),
                Inches(10.333), Inches(5.8),
            )
        except Exception as e:
            logger.warning(f"Chart rendering failed, using text fallback: {e}")
            # Fallback: text-based data display
            self._add_chart_text_fallback(slide, data_labels, data_values)

        self._add_speaker_notes(slide, data.get("notes"))

    def _render_chart_image(
        self,
        chart_type: str,
        labels: list[str],
        values: list[float],
        title: str,
    ) -> BytesIO:
        """Render a chart to a PNG BytesIO using matplotlib."""
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt

        fig, ax = plt.subplots(figsize=(10, 5.5))

        # Convert hex colors to matplotlib format
        colors = [
            f"#{c.red:02x}{c.green:02x}{c.blue:02x}"
            for c in self.CHART_COLORS[:len(labels)]
        ]
        # Pad colors if needed
        while len(colors) < len(labels):
            colors.extend(colors)
        colors = colors[:len(labels)]

        if chart_type == "pie":
            ax.pie(values, labels=labels, colors=colors, autopct="%1.1f%%", startangle=90)
            ax.axis("equal")
        elif chart_type == "horizontal_bar":
            bars = ax.barh(labels, values, color=colors)
            ax.set_xlabel("")
            ax.invert_yaxis()
        elif chart_type == "line":
            ax.plot(labels, values, color=colors[0], marker="o", linewidth=2, markersize=8)
            ax.fill_between(range(len(values)), values, alpha=0.1, color=colors[0])
        else:
            # Default: vertical bar
            bars = ax.bar(labels, values, color=colors)

        ax.set_title("")
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)

        plt.tight_layout()

        buf = BytesIO()
        fig.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor="white")
        plt.close(fig)
        buf.seek(0)
        return buf

    def _add_chart_text_fallback(
        self, slide, labels: list[str], values: list[float]
    ) -> None:
        """Add text-based data representation when chart rendering fails."""
        box = slide.shapes.add_textbox(
            Inches(1), Inches(1.5), Inches(11.333), Inches(5.5)
        )
        tf = box.text_frame
        tf.word_wrap = True

        max_val = max(values) if values else 1
        for i, (label, val) in enumerate(zip(labels, values)):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            bar_len = int((val / max_val) * 20) if max_val > 0 else 0
            bar = "\u2588" * bar_len
            p.text = f"{label}: {bar} {val}"
            p.font.size = Pt(16)
            p.font.color.rgb = self.DARK_TEXT
            p.space_after = Pt(8)

    def _add_closing_slide(self, prs: Presentation, data: dict) -> None:
        """Add a closing/thank you slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_bg(slide, self.SECONDARY_COLOR)

        self._add_text(
            slide, Inches(0.5), Inches(2.8), Inches(12.333), Inches(1.5),
            text=data.get("title", "Thank You"),
            font_size=48, bold=True, color=self.WHITE,
            alignment=PP_ALIGN.CENTER,
        )

        if contact := data.get("contact"):
            self._add_text(
                slide, Inches(0.5), Inches(4.5), Inches(12.333), Inches(0.8),
                text=contact, font_size=20, color=self.WHITE,
                alignment=PP_ALIGN.CENTER,
            )

        self._add_speaker_notes(slide, data.get("notes"))
