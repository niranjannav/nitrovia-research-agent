"""Analytics report renderer.

Renders GeneratedReport objects to PDF and PPTX with embedded matplotlib
chart images. Operates alongside (not replacing) the existing renderers.

Charts are passed as a dict of {chart_name: file_path} pointing to PNG files
in a temporary directory.
"""

import base64
import logging
from io import BytesIO
from pathlib import Path

from app.models.schemas import GeneratedReport

logger = logging.getLogger(__name__)

TEMPLATES_DIR = Path(__file__).parent.parent / "templates"


class AnalyticsPDFRenderer:
    """Renders analytics reports to PDF with embedded chart images."""

    def render(
        self,
        report: GeneratedReport,
        chart_paths: dict[str, str] | None = None,
    ) -> BytesIO:
        """Render an analytics report to PDF with embedded charts.

        Args:
            report: Generated report content
            chart_paths: Dict of chart_name -> PNG file path

        Returns:
            BytesIO containing PDF data
        """
        try:
            from jinja2 import Environment, FileSystemLoader
            from weasyprint import CSS, HTML

            chart_paths = chart_paths or {}

            # Load chart images as base64 for inline embedding
            charts_b64: dict[str, str] = {}
            for name, path in chart_paths.items():
                try:
                    img_bytes = Path(path).read_bytes()
                    charts_b64[name] = base64.b64encode(img_bytes).decode("utf-8")
                except Exception as e:
                    logger.warning(f"Could not load chart '{name}': {e}")

            # Render HTML template
            env = Environment(
                loader=FileSystemLoader(str(TEMPLATES_DIR)),
                autoescape=True,
            )

            # Use analytics template if available, fall back to base template
            try:
                template = env.get_template("analytics_report.html")
            except Exception:
                template = env.get_template("report_base.html")

            html_content = template.render(report=report, charts=charts_b64)

            # Inject chart images into sections if using base template
            if charts_b64 and "<img" not in html_content:
                html_content = _inject_charts_html(html_content, charts_b64)

            css_path = TEMPLATES_DIR / "styles.css"
            css = CSS(filename=str(css_path)) if css_path.exists() else None

            html = HTML(string=html_content, base_url=str(TEMPLATES_DIR))
            pdf_bytes = html.write_pdf(stylesheets=[css] if css else None)

            output = BytesIO(pdf_bytes)
            output.seek(0)
            return output

        except Exception as e:
            logger.error(f"Analytics PDF render failed: {e}")
            raise


class AnalyticsPPTXRenderer:
    """Renders analytics reports to PPTX with chart image slides."""

    PRIMARY_COLOR_HEX = "1E40AF"
    ACCENT_COLOR_HEX = "F59E0B"

    def render(
        self,
        report: GeneratedReport,
        chart_paths: dict[str, str] | None = None,
    ) -> BytesIO:
        """Render an analytics report to PPTX with chart slides.

        Args:
            report: Generated report content
            chart_paths: Dict of chart_name -> PNG file path

        Returns:
            BytesIO containing PPTX data
        """
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt

        chart_paths = chart_paths or {}
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6]  # blank

        # Title slide
        self._add_title_slide(prs, report.title)

        # Executive summary slide
        if report.executive_summary:
            self._add_text_slide(
                prs, "Executive Summary", report.executive_summary[:800], blank_layout
            )

        # Chart slides — one per chart
        CHART_SLIDE_TITLES = {
            "wow_bar": "Week-on-Week Revenue",
            "channel_donut": "Revenue by Channel",
            "channel_comparison": "Channel Performance Comparison",
            "top_products_bar": "Top Products",
            "team_performance": "Team Performance",
            "revenue_trend": "Revenue Trend (24 Months)",
        }

        for chart_name, chart_path in chart_paths.items():
            title = CHART_SLIDE_TITLES.get(chart_name, chart_name.replace("_", " ").title())
            self._add_chart_slide(prs, title, chart_path, blank_layout)

        # Section slides
        for section in report.sections:
            self._add_text_slide(
                prs,
                section.title,
                _truncate(section.content, 600),
                blank_layout,
            )
            for sub in section.subsections or []:
                self._add_text_slide(
                    prs,
                    sub.title,
                    _truncate(sub.content, 600),
                    blank_layout,
                )

        # Recommendations slide
        if report.recommendations:
            recs_text = "\n".join(f"• {r}" for r in report.recommendations[:8])
            self._add_text_slide(prs, "Recommendations", recs_text, blank_layout)

        output = BytesIO()
        prs.save(output)
        output.seek(0)
        return output

    def _add_title_slide(self, prs, title: str) -> None:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(0x1E, 0x40, 0xAF)  # dark blue

        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

    def _add_text_slide(self, prs, title: str, content: str, layout) -> None:
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt

        slide = prs.slides.add_slide(layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

        # Title bar
        title_box = slide.shapes.add_textbox(
            Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.9)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1E, 0x40, 0xAF)

        # Content box
        content_box = slide.shapes.add_textbox(
            Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.8)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(0x37, 0x41, 0x51)

    def _add_chart_slide(self, prs, title: str, chart_path: str, layout) -> None:
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt

        slide = prs.slides.add_slide(layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.7)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1E, 0x40, 0xAF)

        # Chart image
        try:
            slide.shapes.add_picture(
                chart_path,
                Inches(0.5), Inches(1.1),
                width=Inches(12.3), height=Inches(6.0),
            )
        except Exception as e:
            logger.warning(f"Could not add chart image to slide: {e}")


# ── Helpers ───────────────────────────────────────────────────────────────────

def _inject_charts_html(html_content: str, charts_b64: dict[str, str]) -> str:
    """Append chart images before the closing </body> tag."""
    chart_html = "\n<section class='analytics-charts'>\n<h2>Charts</h2>\n"
    for name, b64 in charts_b64.items():
        label = name.replace("_", " ").title()
        chart_html += (
            f"<figure>\n"
            f"  <img src='data:image/png;base64,{b64}' "
            f"  alt='{label}' style='max-width:100%;margin:1em 0;'/>\n"
            f"  <figcaption>{label}</figcaption>\n"
            f"</figure>\n"
        )
    chart_html += "</section>\n"
    return html_content.replace("</body>", chart_html + "</body>")


def _truncate(text: str, max_chars: int) -> str:
    if len(text) <= max_chars:
        return text
    return text[:max_chars].rsplit(" ", 1)[0] + "…"
