"""Document parsing service for multiple file formats.

Each parser extracts content in high-fidelity format and produces
ParsedDocument objects with per-chunk metadata, ready for embedding
and storage in pgvector.
"""

import logging
from abc import ABC, abstractmethod
from io import BytesIO

import fitz  # PyMuPDF
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from app.models.document import (
    DocumentChunk,
    DocumentMetadata,
    ParsedDocument,
    generate_description,
)

logger = logging.getLogger(__name__)

# Maximum characters per chunk for embedding
MAX_CHUNK_CHARS = 4000


class BaseDocumentParser(ABC):
    """Abstract base class for document parsers."""

    @abstractmethod
    def parse(self, file_content: bytes) -> str:
        """Extract text content from document.

        Args:
            file_content: Raw file bytes

        Returns:
            Extracted text content
        """
        pass

    @abstractmethod
    def parse_to_document(
        self,
        file_content: bytes,
        file_name: str,
        user_id: str,
        source_file_id: str | None = None,
    ) -> ParsedDocument:
        """Parse file into a structured ParsedDocument with metadata.

        Args:
            file_content: Raw file bytes
            file_name: Original filename
            user_id: Owner identifier for search filtering
            source_file_id: Optional reference to source_files table

        Returns:
            ParsedDocument with chunks and metadata
        """
        pass

    @abstractmethod
    def supported_extensions(self) -> list[str]:
        """Return list of supported file extensions."""
        pass

    def _get_file_type(self) -> str:
        """Get the file type string for metadata."""
        exts = self.supported_extensions()
        return exts[0].lstrip(".") if exts else "unknown"

    def _create_metadata(
        self,
        file_name: str,
        user_id: str,
        description: str,
        source_file_id: str | None = None,
        chunk_index: int = 0,
        total_chunks: int = 1,
        page_number: int | None = None,
        sheet_name: str | None = None,
    ) -> DocumentMetadata:
        """Create metadata for a document chunk."""
        return DocumentMetadata(
            file_type=self._get_file_type(),
            file_name=file_name,
            description=description,
            user_id=user_id,
            source_file_id=source_file_id,
            chunk_index=chunk_index,
            total_chunks=total_chunks,
            page_number=page_number,
            sheet_name=sheet_name,
        )


class PDFParser(BaseDocumentParser):
    """Parser for PDF files using PyMuPDF."""

    def parse(self, file_content: bytes) -> str:
        """Extract text from PDF."""
        try:
            doc = fitz.open(stream=file_content, filetype="pdf")
            text_parts = []

            for page_num, page in enumerate(doc, 1):
                page_text = page.get_text()
                if page_text.strip():
                    text_parts.append(f"[Page {page_num}]\n{page_text}")

            doc.close()
            return "\n\n".join(text_parts)

        except Exception as e:
            logger.error(f"PDF parsing error: {e}")
            raise ValueError(f"Failed to parse PDF: {str(e)}")

    def parse_to_document(
        self,
        file_content: bytes,
        file_name: str,
        user_id: str,
        source_file_id: str | None = None,
    ) -> ParsedDocument:
        """Parse PDF into structured document with per-page chunks."""
        try:
            doc = fitz.open(stream=file_content, filetype="pdf")
            chunks: list[DocumentChunk] = []
            text_parts: list[str] = []

            page_count = len(doc)
            for page_num, page in enumerate(doc, 1):
                page_text = page.get_text()
                if page_text.strip():
                    text_parts.append(f"[Page {page_num}]\n{page_text}")

            raw_text = "\n\n".join(text_parts)
            description = generate_description(raw_text)

            # Create chunks per page
            for page_num, page in enumerate(doc, 1):
                page_text = page.get_text()
                if page_text.strip():
                    metadata = self._create_metadata(
                        file_name=file_name,
                        user_id=user_id,
                        description=description,
                        source_file_id=source_file_id,
                        chunk_index=len(chunks),
                        total_chunks=page_count,
                        page_number=page_num,
                    )
                    chunks.append(DocumentChunk(content=page_text.strip(), metadata=metadata))

            doc.close()

            # Update total_chunks in all metadata
            for chunk in chunks:
                chunk.metadata.total_chunks = len(chunks)

            return ParsedDocument(
                file_name=file_name,
                file_type="pdf",
                chunks=chunks,
                raw_text=raw_text,
                description=description,
            )

        except Exception as e:
            logger.error(f"PDF parsing error: {e}")
            raise ValueError(f"Failed to parse PDF: {str(e)}")

    def supported_extensions(self) -> list[str]:
        return [".pdf"]


class DOCXParser(BaseDocumentParser):
    """Parser for DOCX files using python-docx."""

    def parse(self, file_content: bytes) -> str:
        """Extract text from DOCX."""
        try:
            doc = Document(BytesIO(file_content))
            text_parts = []

            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)

            # Also extract text from tables
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_text:
                        table_text.append(row_text)
                if table_text:
                    text_parts.append("\n[Table]\n" + "\n".join(table_text))

            return "\n\n".join(text_parts)

        except Exception as e:
            logger.error(f"DOCX parsing error: {e}")
            raise ValueError(f"Failed to parse DOCX: {str(e)}")

    def parse_to_document(
        self,
        file_content: bytes,
        file_name: str,
        user_id: str,
        source_file_id: str | None = None,
    ) -> ParsedDocument:
        """Parse DOCX into structured document with logical chunks."""
        try:
            doc = Document(BytesIO(file_content))
            chunks: list[DocumentChunk] = []
            raw_text = self.parse(file_content)
            description = generate_description(raw_text)

            # Build chunks from paragraphs, grouping by size
            current_chunk_text: list[str] = []
            current_chars = 0

            for para in doc.paragraphs:
                if para.text.strip():
                    # Start a new chunk if adding this paragraph would exceed limit
                    if current_chars + len(para.text) > MAX_CHUNK_CHARS and current_chunk_text:
                        content = "\n\n".join(current_chunk_text)
                        metadata = self._create_metadata(
                            file_name=file_name,
                            user_id=user_id,
                            description=description,
                            source_file_id=source_file_id,
                            chunk_index=len(chunks),
                        )
                        chunks.append(DocumentChunk(content=content, metadata=metadata))
                        current_chunk_text = []
                        current_chars = 0

                    current_chunk_text.append(para.text)
                    current_chars += len(para.text)

            # Add tables as separate chunks
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_text:
                        table_text.append(row_text)
                if table_text:
                    content = "[Table]\n" + "\n".join(table_text)
                    metadata = self._create_metadata(
                        file_name=file_name,
                        user_id=user_id,
                        description=description,
                        source_file_id=source_file_id,
                        chunk_index=len(chunks),
                    )
                    chunks.append(DocumentChunk(content=content, metadata=metadata))

            # Flush remaining paragraph text
            if current_chunk_text:
                content = "\n\n".join(current_chunk_text)
                metadata = self._create_metadata(
                    file_name=file_name,
                    user_id=user_id,
                    description=description,
                    source_file_id=source_file_id,
                    chunk_index=len(chunks),
                )
                chunks.append(DocumentChunk(content=content, metadata=metadata))

            # If no chunks were created, create one from raw text
            if not chunks and raw_text.strip():
                metadata = self._create_metadata(
                    file_name=file_name,
                    user_id=user_id,
                    description=description,
                    source_file_id=source_file_id,
                )
                chunks.append(DocumentChunk(content=raw_text.strip(), metadata=metadata))

            # Update total_chunks
            for chunk in chunks:
                chunk.metadata.total_chunks = len(chunks)

            return ParsedDocument(
                file_name=file_name,
                file_type="docx",
                chunks=chunks,
                raw_text=raw_text,
                description=description,
            )

        except Exception as e:
            logger.error(f"DOCX parsing error: {e}")
            raise ValueError(f"Failed to parse DOCX: {str(e)}")

    def supported_extensions(self) -> list[str]:
        return [".docx", ".doc"]


class XLSXParser(BaseDocumentParser):
    """Parser for Excel files using openpyxl.

    Creates document chunks per sheet tab, splitting large sheets into
    multiple chunks to keep each within MAX_CHUNK_CHARS. The parse()
    method produces a compact summary (headers + sample rows + stats)
    for context building, while parse_to_document() produces properly
    chunked content for embedding.
    """

    # Max rows to scan for statistics/type inference
    MAX_STAT_ROWS = 100
    # Max data rows to include in sample output (parse() summary)
    MAX_SAMPLE_ROWS = 30
    # Max characters per sheet in the parse() summary
    MAX_CHARS_PER_SHEET = 15_000
    # Max total characters for the entire parse() output
    MAX_TOTAL_CHARS = 80_000
    # Max data rows per sheet in parse_to_document() chunks
    MAX_DATA_ROWS = 500

    def _extract_headers(
        self,
        sheet,
        rows_count: int,
    ) -> tuple[list[str], int]:
        """Extract header row from a sheet.

        Returns:
            Tuple of (header values, header row index 1-based)
        """
        header_row_idx = 1
        for row in sheet.iter_rows(min_row=1, max_row=min(5, rows_count), values_only=True):
            row_vals = [str(c).strip() if c is not None else "" for c in row]
            if any(v for v in row_vals):
                return row_vals, header_row_idx
            header_row_idx += 1
        return [], header_row_idx

    def _infer_column_types(
        self,
        sheet,
        cols_count: int,
        data_start: int,
        rows_count: int,
    ) -> tuple[list[str], list[list]]:
        """Infer column data types from sample data rows.

        Returns:
            Tuple of (column types list, column values list)
        """
        col_types: list[str] = ["unknown"] * cols_count
        col_values: list[list] = [[] for _ in range(cols_count)]

        scan_end = min(data_start + self.MAX_STAT_ROWS, rows_count + 1)
        for row in sheet.iter_rows(min_row=data_start, max_row=scan_end, values_only=True):
            for ci, cell in enumerate(row):
                if ci >= cols_count:
                    break
                if cell is not None:
                    col_values[ci].append(cell)

        for ci in range(cols_count):
            vals = col_values[ci]
            if not vals:
                col_types[ci] = "empty"
                continue
            num_count = sum(1 for v in vals if isinstance(v, (int, float)))
            str_count = sum(1 for v in vals if isinstance(v, str))
            if num_count > len(vals) * 0.7:
                col_types[ci] = "numeric"
            elif str_count > len(vals) * 0.7:
                col_types[ci] = "text"
            else:
                col_types[ci] = "mixed"

        return col_types, col_values

    def _format_row(self, row) -> str:
        """Format a row tuple into a pipe-separated string."""
        return " | ".join(str(cell) if cell is not None else "" for cell in row)

    def parse(self, file_content: bytes) -> str:
        """Extract structured summary text from Excel spreadsheet.

        Produces a compact representation per sheet: headers with types,
        sample data rows (capped), column statistics, formulas, and
        completeness warnings. Output is capped at MAX_TOTAL_CHARS.
        """
        from statistics import mean, median, stdev

        try:
            wb = load_workbook(BytesIO(file_content), data_only=True)
            try:
                wb_formulas = load_workbook(BytesIO(file_content), data_only=False)
            except Exception:
                wb_formulas = None

            text_parts = []
            total_chars = 0

            for sheet_idx, sheet in enumerate(wb.worksheets):
                rows_count = sheet.max_row or 0
                cols_count = sheet.max_column or 0

                if rows_count == 0 or cols_count == 0:
                    continue

                section = [f"[Sheet: {sheet.title} | {rows_count} rows × {cols_count} cols]"]

                headers, header_row_idx = self._extract_headers(sheet, rows_count)
                data_start = header_row_idx + 1

                col_types, col_values = self._infer_column_types(
                    sheet, cols_count, data_start, rows_count
                )

                # --- Headers with types ---
                if headers:
                    header_line = " | ".join(
                        f"{h} ({col_types[i]})" if i < len(col_types) else h
                        for i, h in enumerate(headers)
                    )
                    section.append(f"Headers: {header_line}")

                # --- Data sample (capped) ---
                section.append("")
                section.append("Data:")
                sample_end = min(data_start + self.MAX_SAMPLE_ROWS, rows_count + 1)
                sheet_data_chars = 0
                for row in sheet.iter_rows(min_row=data_start, max_row=sample_end, values_only=True):
                    row_text = self._format_row(row)
                    if row_text.strip(" |"):
                        section.append(row_text)
                        sheet_data_chars += len(row_text)
                        if sheet_data_chars > self.MAX_CHARS_PER_SHEET:
                            section.append("... (sheet output truncated)")
                            break

                remaining = rows_count - (sample_end - data_start) - data_start + 1
                if remaining > 0:
                    section.append(f"... ({remaining} more rows)")

                # --- Per-column statistics for numeric columns ---
                stats_lines = []
                for ci in range(cols_count):
                    if col_types[ci] != "numeric":
                        continue
                    nums = [v for v in col_values[ci] if isinstance(v, (int, float))]
                    if len(nums) < 2:
                        continue
                    header_name = headers[ci] if ci < len(headers) and headers[ci] else f"Col{ci+1}"
                    try:
                        stats_lines.append(
                            f"  {header_name}: "
                            f"min={min(nums):.2f}, max={max(nums):.2f}, "
                            f"mean={mean(nums):.2f}, median={median(nums):.2f}, "
                            f"sum={sum(nums):.2f}, stdev={stdev(nums):.2f}, "
                            f"count={len(nums)}"
                        )
                    except Exception:
                        pass

                if stats_lines:
                    section.append("")
                    section.append("Column Statistics:")
                    section.extend(stats_lines)

                # --- Formula detection ---
                if wb_formulas and sheet_idx < len(wb_formulas.worksheets):
                    formula_sheet = wb_formulas.worksheets[sheet_idx]
                    formulas_found = []
                    for row in formula_sheet.iter_rows(
                        min_row=1,
                        max_row=min(rows_count, 200),
                        max_col=cols_count,
                    ):
                        for cell in row:
                            if cell.value and isinstance(cell.value, str) and cell.value.startswith("="):
                                formulas_found.append(f"  {cell.coordinate}: {cell.value}")
                                if len(formulas_found) >= 20:
                                    break
                        if len(formulas_found) >= 20:
                            break

                    if formulas_found:
                        section.append("")
                        section.append("Formulas detected:")
                        section.extend(formulas_found)

                # --- Empty-cell analysis ---
                scan_end_stat = min(data_start + self.MAX_STAT_ROWS, rows_count + 1)
                total_data_rows = min(scan_end_stat - data_start, rows_count)
                if total_data_rows > 0:
                    empty_lines = []
                    for ci in range(cols_count):
                        empty_count = total_data_rows - len(col_values[ci])
                        if empty_count > 0:
                            pct = (empty_count / total_data_rows) * 100
                            if pct > 10:
                                header_name = headers[ci] if ci < len(headers) and headers[ci] else f"Col{ci+1}"
                                empty_lines.append(f"  {header_name}: {pct:.0f}% empty")

                    if empty_lines:
                        section.append("")
                        section.append("Data completeness warnings:")
                        section.extend(empty_lines)

                sheet_text = "\n".join(section)
                text_parts.append(sheet_text)
                total_chars += len(sheet_text)

                # Stop processing more sheets if we've hit the overall cap
                if total_chars >= self.MAX_TOTAL_CHARS:
                    remaining_sheets = len(wb.worksheets) - sheet_idx - 1
                    if remaining_sheets > 0:
                        text_parts.append(
                            f"... ({remaining_sheets} more sheet(s) omitted, "
                            f"output capped at {self.MAX_TOTAL_CHARS} chars)"
                        )
                    break

            wb.close()
            if wb_formulas:
                wb_formulas.close()

            return "\n\n".join(text_parts)

        except Exception as e:
            logger.error(f"XLSX parsing error: {e}")
            raise ValueError(f"Failed to parse Excel file: {str(e)}")

    def parse_to_document(
        self,
        file_content: bytes,
        file_name: str,
        user_id: str,
        source_file_id: str | None = None,
    ) -> ParsedDocument:
        """Parse Excel into structured document with properly-sized chunks.

        Each sheet is split into one or more chunks, each capped at
        MAX_CHUNK_CHARS. The raw_text is a compact summary (from parse())
        suitable for context building without token explosion.
        """
        try:
            wb = load_workbook(BytesIO(file_content), data_only=True)
            chunks: list[DocumentChunk] = []

            # Generate compact summary for raw_text (used by build_context)
            raw_text = self.parse(file_content)
            description = generate_description(raw_text)

            for sheet in wb.worksheets:
                rows_count = sheet.max_row or 0
                cols_count = sheet.max_column or 0
                if rows_count == 0 or cols_count == 0:
                    continue

                headers, header_row_idx = self._extract_headers(sheet, rows_count)
                data_start = header_row_idx + 1
                header_text = self._format_row(headers) if headers else ""

                sheet_prefix = f"[Sheet: {sheet.title}]\n"
                if header_text:
                    sheet_prefix += f"Headers: {header_text}\n"

                # Iterate data rows, splitting into chunks at MAX_CHUNK_CHARS
                current_lines: list[str] = []
                current_chars = len(sheet_prefix)
                rows_emitted = 0
                max_rows = min(
                    rows_count - data_start + 1,
                    self.MAX_DATA_ROWS,
                )

                for row in sheet.iter_rows(
                    min_row=data_start,
                    max_row=data_start + max_rows - 1,
                    values_only=True,
                ):
                    row_text = self._format_row(row)
                    if not row_text.strip(" |"):
                        continue

                    row_len = len(row_text) + 1  # +1 for newline
                    # If adding this row exceeds the chunk limit, flush
                    if current_chars + row_len > MAX_CHUNK_CHARS and current_lines:
                        content = sheet_prefix + "\n".join(current_lines)
                        metadata = self._create_metadata(
                            file_name=file_name,
                            user_id=user_id,
                            description=description,
                            source_file_id=source_file_id,
                            chunk_index=len(chunks),
                            sheet_name=sheet.title,
                        )
                        chunks.append(DocumentChunk(content=content, metadata=metadata))
                        current_lines = []
                        current_chars = len(sheet_prefix)

                    current_lines.append(row_text)
                    current_chars += row_len
                    rows_emitted += 1

                # Flush remaining rows for this sheet
                if current_lines:
                    remaining = rows_count - data_start + 1 - rows_emitted
                    if remaining > 0:
                        current_lines.append(f"... ({remaining} more rows)")
                    content = sheet_prefix + "\n".join(current_lines)
                    metadata = self._create_metadata(
                        file_name=file_name,
                        user_id=user_id,
                        description=description,
                        source_file_id=source_file_id,
                        chunk_index=len(chunks),
                        sheet_name=sheet.title,
                    )
                    chunks.append(DocumentChunk(content=content, metadata=metadata))

            wb.close()

            # If no chunks, create one from raw text
            if not chunks and raw_text.strip():
                metadata = self._create_metadata(
                    file_name=file_name,
                    user_id=user_id,
                    description=description,
                    source_file_id=source_file_id,
                )
                chunks.append(DocumentChunk(content=raw_text.strip(), metadata=metadata))

            # Update total_chunks
            for chunk in chunks:
                chunk.metadata.total_chunks = len(chunks)

            logger.info(
                f"XLSX parsed: {file_name} -> {len(chunks)} chunks, "
                f"raw_text={len(raw_text)} chars"
            )

            return ParsedDocument(
                file_name=file_name,
                file_type="xlsx",
                chunks=chunks,
                raw_text=raw_text,
                description=description,
            )

        except Exception as e:
            logger.error(f"XLSX parsing error: {e}")
            raise ValueError(f"Failed to parse Excel file: {str(e)}")

    def supported_extensions(self) -> list[str]:
        return [".xlsx", ".xls"]


class PPTXParser(BaseDocumentParser):
    """Parser for PowerPoint files using python-pptx."""

    def parse(self, file_content: bytes) -> str:
        """Extract structured text from PowerPoint presentation."""
        try:
            prs = Presentation(BytesIO(file_content))
            text_parts = []
            total_slides = len(prs.slides)

            text_parts.append(f"[Presentation: {total_slides} slides]")

            for i, slide in enumerate(prs.slides, 1):
                slide_text = []

                # Detect layout type
                layout_name = "unknown"
                try:
                    if slide.slide_layout and slide.slide_layout.name:
                        layout_name = slide.slide_layout.name
                except Exception:
                    pass

                slide_text.append(f"[Slide {i}/{total_slides} | Layout: {layout_name}]")

                # Track content types found
                has_images = False
                has_charts = False
                has_tables = False

                for shape in slide.shapes:
                    # Text content
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)

                    # Table extraction
                    if shape.has_table:
                        has_tables = True
                        table = shape.table
                        table_rows = []
                        for row in table.rows:
                            row_text = " | ".join(
                                cell.text.strip()
                                for cell in row.cells
                                if cell.text.strip()
                            )
                            if row_text:
                                table_rows.append(row_text)
                        if table_rows:
                            slide_text.append("[Table]")
                            slide_text.extend(table_rows)

                    # Image detection
                    if shape.shape_type and shape.shape_type == 13:  # Picture
                        has_images = True

                    # Chart detection
                    if hasattr(shape, "has_chart") and shape.has_chart:
                        has_charts = True

                # Media indicators
                media_notes = []
                if has_images:
                    media_notes.append("contains image(s)")
                if has_charts:
                    media_notes.append("contains chart(s)")
                if has_tables:
                    media_notes.append("contains table(s)")
                if media_notes:
                    slide_text.append(f"[Media: {', '.join(media_notes)}]")

                # Speaker notes
                try:
                    if slide.has_notes_slide and slide.notes_slide:
                        notes_text = slide.notes_slide.notes_text_frame.text.strip()
                        if notes_text:
                            slide_text.append(f"[Speaker Notes: {notes_text}]")
                except Exception:
                    pass

                if len(slide_text) > 1:  # More than just the slide marker
                    text_parts.append("\n".join(slide_text))

            return "\n\n".join(text_parts)

        except Exception as e:
            logger.error(f"PPTX parsing error: {e}")
            raise ValueError(f"Failed to parse PowerPoint file: {str(e)}")

    def parse_to_document(
        self,
        file_content: bytes,
        file_name: str,
        user_id: str,
        source_file_id: str | None = None,
    ) -> ParsedDocument:
        """Parse PowerPoint into structured document with per-slide chunks."""
        try:
            prs = Presentation(BytesIO(file_content))
            chunks: list[DocumentChunk] = []
            raw_text = self.parse(file_content)
            description = generate_description(raw_text)

            for i, slide in enumerate(prs.slides, 1):
                slide_content_parts: list[str] = []

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content_parts.append(shape.text)

                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = " | ".join(
                                cell.text.strip()
                                for cell in row.cells
                                if cell.text.strip()
                            )
                            if row_text:
                                slide_content_parts.append(row_text)

                if slide_content_parts:
                    content = f"[Slide {i}]\n" + "\n".join(slide_content_parts)
                    metadata = self._create_metadata(
                        file_name=file_name,
                        user_id=user_id,
                        description=description,
                        source_file_id=source_file_id,
                        chunk_index=len(chunks),
                        page_number=i,
                    )
                    chunks.append(DocumentChunk(content=content, metadata=metadata))

            # If no chunks, create one from raw text
            if not chunks and raw_text.strip():
                metadata = self._create_metadata(
                    file_name=file_name,
                    user_id=user_id,
                    description=description,
                    source_file_id=source_file_id,
                )
                chunks.append(DocumentChunk(content=raw_text.strip(), metadata=metadata))

            # Update total_chunks
            for chunk in chunks:
                chunk.metadata.total_chunks = len(chunks)

            return ParsedDocument(
                file_name=file_name,
                file_type="pptx",
                chunks=chunks,
                raw_text=raw_text,
                description=description,
            )

        except Exception as e:
            logger.error(f"PPTX parsing error: {e}")
            raise ValueError(f"Failed to parse PowerPoint file: {str(e)}")

    def supported_extensions(self) -> list[str]:
        return [".pptx", ".ppt"]


class ParserFactory:
    """Factory for getting document parsers by file type."""

    _parsers: dict[str, BaseDocumentParser] = {}

    @classmethod
    def register(cls, parser: BaseDocumentParser) -> None:
        """Register a parser for its supported extensions."""
        for ext in parser.supported_extensions():
            cls._parsers[ext.lower()] = parser

    @classmethod
    def get_parser(cls, file_extension: str) -> BaseDocumentParser:
        """Get parser for file extension.

        Args:
            file_extension: File extension (e.g., ".pdf")

        Returns:
            BaseDocumentParser instance

        Raises:
            ValueError: If no parser exists for extension
        """
        ext = file_extension.lower()
        if not ext.startswith("."):
            ext = f".{ext}"

        if ext not in cls._parsers:
            raise ValueError(f"Unsupported file type: {ext}")

        return cls._parsers[ext]

    @classmethod
    def parse_file(cls, file_content: bytes, file_extension: str) -> str:
        """Parse file content using appropriate parser.

        Args:
            file_content: Raw file bytes
            file_extension: File extension

        Returns:
            Extracted text content
        """
        parser = cls.get_parser(file_extension)
        return parser.parse(file_content)

    @classmethod
    def parse_file_to_document(
        cls,
        file_content: bytes,
        file_extension: str,
        file_name: str,
        user_id: str,
        source_file_id: str | None = None,
    ) -> ParsedDocument:
        """Parse file content into a structured ParsedDocument.

        Args:
            file_content: Raw file bytes
            file_extension: File extension
            file_name: Original filename
            user_id: Owner identifier for search filtering
            source_file_id: Optional reference to source_files table

        Returns:
            ParsedDocument with chunks and metadata
        """
        parser = cls.get_parser(file_extension)
        return parser.parse_to_document(
            file_content=file_content,
            file_name=file_name,
            user_id=user_id,
            source_file_id=source_file_id,
        )


# Register all parsers
ParserFactory.register(PDFParser())
ParserFactory.register(DOCXParser())
ParserFactory.register(XLSXParser())
ParserFactory.register(PPTXParser())
