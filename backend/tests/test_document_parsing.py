"""Tests for document parsing with metadata and document models.

Tests the enhanced document parsers that produce ParsedDocument objects
with metadata, the document model utilities, and the research planner.

Run with: pytest tests/test_document_parsing.py -v
"""

import io

import pytest
from docx import Document as DocxDocument
from openpyxl import Workbook
from pptx import Presentation

from app.models.document import (
    DocumentChunk,
    DocumentMetadata,
    ParsedDocument,
    generate_description,
)
from app.services.document_parser import (
    DOCXParser,
    ParserFactory,
    PPTXParser,
    XLSXParser,
)

# ============================================
# Document Model Tests
# ============================================


class TestDocumentMetadata:
    """Test DocumentMetadata model."""

    def test_metadata_creation(self):
        """Test creating metadata with required fields."""
        meta = DocumentMetadata(
            file_type="pdf",
            file_name="test.pdf",
            description="A test document",
            user_id="user-123",
        )
        assert meta.file_type == "pdf"
        assert meta.file_name == "test.pdf"
        assert meta.description == "A test document"
        assert meta.user_id == "user-123"
        assert meta.chunk_index == 0
        assert meta.total_chunks == 1

    def test_metadata_to_dict(self):
        """Test metadata serialization."""
        meta = DocumentMetadata(
            file_type="xlsx",
            file_name="data.xlsx",
            description="Excel data file",
            user_id="user-456",
            source_file_id="file-789",
            chunk_index=2,
            total_chunks=5,
            sheet_name="Sheet1",
        )
        d = meta.to_dict()
        assert d["file_type"] == "xlsx"
        assert d["file_name"] == "data.xlsx"
        assert d["user_id"] == "user-456"
        assert d["source_file_id"] == "file-789"
        assert d["chunk_index"] == 2
        assert d["total_chunks"] == 5
        assert d["sheet_name"] == "Sheet1"

    def test_metadata_optional_fields_excluded(self):
        """Test that None optional fields are excluded from dict."""
        meta = DocumentMetadata(
            file_type="pdf",
            file_name="test.pdf",
            description="Test",
            user_id="user-1",
        )
        d = meta.to_dict()
        assert "page_number" not in d
        assert "sheet_name" not in d


class TestDocumentChunk:
    """Test DocumentChunk model."""

    def test_chunk_creation(self):
        """Test creating a document chunk."""
        meta = DocumentMetadata(
            file_type="pdf",
            file_name="test.pdf",
            description="Test",
            user_id="user-1",
        )
        chunk = DocumentChunk(content="Hello world", metadata=meta)
        assert chunk.content == "Hello world"
        assert chunk.metadata.file_type == "pdf"

    def test_chunk_to_dict(self):
        """Test chunk serialization."""
        meta = DocumentMetadata(
            file_type="docx",
            file_name="doc.docx",
            description="A document",
            user_id="user-1",
        )
        chunk = DocumentChunk(content="Test content", metadata=meta)
        d = chunk.to_dict()
        assert d["content"] == "Test content"
        assert d["metadata"]["file_type"] == "docx"


class TestParsedDocument:
    """Test ParsedDocument model."""

    def test_parsed_document_creation(self):
        """Test creating a parsed document."""
        meta = DocumentMetadata(
            file_type="pdf",
            file_name="test.pdf",
            description="Test",
            user_id="user-1",
        )
        chunks = [
            DocumentChunk(content="Page 1 content", metadata=meta),
            DocumentChunk(content="Page 2 content", metadata=meta),
        ]
        doc = ParsedDocument(
            file_name="test.pdf",
            file_type="pdf",
            chunks=chunks,
            raw_text="Page 1 content\nPage 2 content",
            description="Test",
        )
        assert doc.total_chunks == 2
        assert doc.file_name == "test.pdf"

    def test_parsed_document_to_dict(self):
        """Test parsed document serialization."""
        meta = DocumentMetadata(
            file_type="pdf",
            file_name="test.pdf",
            description="Test",
            user_id="user-1",
        )
        chunks = [DocumentChunk(content="Content", metadata=meta)]
        doc = ParsedDocument(
            file_name="test.pdf",
            file_type="pdf",
            chunks=chunks,
            raw_text="Content",
            description="Test",
        )
        d = doc.to_dict()
        assert d["file_name"] == "test.pdf"
        assert d["total_chunks"] == 1
        assert len(d["chunks"]) == 1


class TestGenerateDescription:
    """Test the generate_description utility."""

    def test_basic_description(self):
        """Test generating description from content."""
        content = "This is a research report about climate change impacts on agriculture."
        desc = generate_description(content)
        assert "research report" in desc.lower() or "climate" in desc.lower()

    def test_empty_content(self):
        """Test description for empty content."""
        assert generate_description("") == "Empty document"
        assert generate_description("   ") == "Empty document"

    def test_description_truncation(self):
        """Test description is truncated to max_length."""
        content = "A" * 1000
        desc = generate_description(content, max_length=100)
        assert len(desc) <= 103  # max_length + "..."

    def test_skips_structural_markers(self):
        """Test that structural markers are skipped."""
        content = "[Page 1]\n[Sheet: Data]\nActual content here"
        desc = generate_description(content)
        assert desc.startswith("Actual content")

    def test_multiline_content(self):
        """Test with multiple lines."""
        content = "First line of the document.\nSecond line with details.\nThird line."
        desc = generate_description(content)
        assert "First line" in desc


# ============================================
# Parser Tests (file-type specific)
# ============================================


def _create_docx_bytes(paragraphs: list[str]) -> bytes:
    """Helper: create a DOCX file in memory."""
    doc = DocxDocument()
    for para in paragraphs:
        doc.add_paragraph(para)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _create_xlsx_bytes(sheets: dict[str, list[list[str]]]) -> bytes:
    """Helper: create an XLSX file in memory."""
    wb = Workbook()
    first = True
    for sheet_name, rows in sheets.items():
        if first:
            ws = wb.active
            ws.title = sheet_name
            first = False
        else:
            ws = wb.create_sheet(sheet_name)
        for row in rows:
            ws.append(row)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _create_pptx_bytes(slides: list[str]) -> bytes:
    """Helper: create a PPTX file in memory."""
    prs = Presentation()
    for slide_text in slides:
        slide_layout = prs.slide_layouts[1]  # Title + Content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = slide_text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


class TestDOCXParserWithMetadata:
    """Test DOCX parser producing ParsedDocument."""

    def test_parse_to_document(self):
        """Test parsing DOCX to structured document."""
        content = _create_docx_bytes(["First paragraph.", "Second paragraph."])
        parser = DOCXParser()
        doc = parser.parse_to_document(
            file_content=content,
            file_name="test.docx",
            user_id="user-123",
            source_file_id="file-456",
        )

        assert doc.file_name == "test.docx"
        assert doc.file_type == "docx"
        assert len(doc.chunks) >= 1
        assert doc.raw_text  # Should have raw text
        assert doc.description  # Should have auto-generated description

        # Check metadata on chunks
        for chunk in doc.chunks:
            assert chunk.metadata.file_type == "docx"
            assert chunk.metadata.file_name == "test.docx"
            assert chunk.metadata.user_id == "user-123"
            assert chunk.metadata.source_file_id == "file-456"
            assert chunk.metadata.description

    def test_parse_backward_compat(self):
        """Test that raw parse() still works."""
        content = _create_docx_bytes(["Hello world."])
        parser = DOCXParser()
        text = parser.parse(content)
        assert "Hello world" in text


class TestXLSXParserWithMetadata:
    """Test XLSX parser producing ParsedDocument."""

    def test_parse_to_document_per_sheet(self):
        """Test parsing XLSX creates chunk per sheet."""
        content = _create_xlsx_bytes({
            "Sales": [["Product", "Revenue"], ["Widget", "1000"]],
            "Costs": [["Item", "Amount"], ["Materials", "500"]],
        })
        parser = XLSXParser()
        doc = parser.parse_to_document(
            file_content=content,
            file_name="data.xlsx",
            user_id="user-123",
        )

        assert doc.file_name == "data.xlsx"
        assert doc.file_type == "xlsx"
        assert len(doc.chunks) == 2  # One chunk per sheet

        # Verify sheet names in metadata
        sheet_names = {chunk.metadata.sheet_name for chunk in doc.chunks}
        assert "Sales" in sheet_names
        assert "Costs" in sheet_names

        # Verify all chunks have correct metadata
        for chunk in doc.chunks:
            assert chunk.metadata.file_type == "xlsx"
            assert chunk.metadata.user_id == "user-123"

    def test_parse_backward_compat(self):
        """Test that raw parse() still works."""
        content = _create_xlsx_bytes({
            "Sheet1": [["A", "B"], ["1", "2"]],
        })
        parser = XLSXParser()
        text = parser.parse(content)
        assert "Sheet1" in text


class TestPPTXParserWithMetadata:
    """Test PPTX parser producing ParsedDocument."""

    def test_parse_to_document_per_slide(self):
        """Test parsing PPTX creates chunk per slide."""
        content = _create_pptx_bytes(["Introduction", "Main Content", "Conclusion"])
        parser = PPTXParser()
        doc = parser.parse_to_document(
            file_content=content,
            file_name="presentation.pptx",
            user_id="user-123",
        )

        assert doc.file_name == "presentation.pptx"
        assert doc.file_type == "pptx"
        assert len(doc.chunks) == 3  # One chunk per slide

        # Verify page numbers
        page_numbers = [chunk.metadata.page_number for chunk in doc.chunks]
        assert 1 in page_numbers
        assert 2 in page_numbers
        assert 3 in page_numbers

    def test_parse_backward_compat(self):
        """Test that raw parse() still works."""
        content = _create_pptx_bytes(["Test Slide"])
        parser = PPTXParser()
        text = parser.parse(content)
        assert "Test Slide" in text


class TestParserFactory:
    """Test ParserFactory with enhanced parsing."""

    def test_parse_file_to_document(self):
        """Test factory method for structured parsing."""
        content = _create_docx_bytes(["Test content for factory."])
        doc = ParserFactory.parse_file_to_document(
            file_content=content,
            file_extension=".docx",
            file_name="factory_test.docx",
            user_id="user-789",
        )

        assert isinstance(doc, ParsedDocument)
        assert doc.file_name == "factory_test.docx"
        assert doc.chunks
        assert doc.chunks[0].metadata.user_id == "user-789"

    def test_parse_file_backward_compat(self):
        """Test that original parse_file still works."""
        content = _create_docx_bytes(["Backward compatible."])
        text = ParserFactory.parse_file(content, ".docx")
        assert "Backward compatible" in text

    def test_unsupported_extension(self):
        """Test error for unsupported file type."""
        with pytest.raises(ValueError, match="Unsupported file type"):
            ParserFactory.get_parser(".xyz")


class TestResearchPlanModel:
    """Test ResearchPlan model validation."""

    def test_research_plan_creation(self):
        """Test creating a research plan."""
        from app.services.research_planner import ResearchPlan

        plan = ResearchPlan(
            title="Test Report",
            summary="Research plan for testing",
            questions=[
                "What are the main findings?",
                "What data supports the conclusions?",
                "What recommendations can be made?",
            ],
        )
        assert plan.title == "Test Report"
        assert len(plan.questions) == 3

    def test_research_plan_min_questions(self):
        """Test that plan requires minimum questions."""
        from app.services.research_planner import ResearchPlan

        with pytest.raises(Exception):
            ResearchPlan(
                title="Test",
                summary="Test",
                questions=["Only one question"],
            )


class TestResearchPlannerFallback:
    """Test ResearchPlanner fallback (no LLM required)."""

    def test_fallback_plan(self):
        """Test fallback plan generation without LLM."""
        from app.services.research_planner import ResearchPlanner

        # Create planner but use fallback directly
        planner = ResearchPlanner.__new__(ResearchPlanner)
        plan = planner._fallback_plan("Climate Change Impact")

        assert plan.title == "Climate Change Impact"
        assert len(plan.questions) >= 5
        assert any("Climate Change Impact" in q for q in plan.questions)

    def test_fallback_plan_with_instructions(self):
        """Test fallback plan with custom instructions."""
        from app.services.research_planner import ResearchPlanner

        planner = ResearchPlanner.__new__(ResearchPlanner)
        plan = planner._fallback_plan(
            "Market Analysis",
            custom_instructions="Focus on emerging markets in Asia",
        )

        assert len(plan.questions) >= 6  # 5 base + 1 custom
        assert "Focus on emerging markets in Asia" in plan.questions
